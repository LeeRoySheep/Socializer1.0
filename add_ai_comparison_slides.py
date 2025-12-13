#!/usr/bin/env python3
"""
Add AI Comparison Results to PowerPoint Presentation
====================================================

PURPOSE:
    Adds comprehensive AI provider comparison slides to the existing presentation
    after videos and screenshots.

LOCATION:
    add_ai_comparison_slides.py

DEPENDENCIES:
    - python-pptx: PowerPoint manipulation
    - Socializer_Presentation.pptx: Existing presentation
    - AI_PROVIDER_COMPARISON.md: Source data

SLIDES ADDED:
    1. AI Provider Overview (comparison table)
    2. Speed Comparison (visual chart)
    3. Cost Comparison (breakdown)
    4. Quality & Accuracy Analysis
    5. Recommendations

USAGE:
    python add_ai_comparison_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load existing presentation
print("📂 Loading presentation...")
prs = Presentation('Socializer_Presentation.pptx')
print(f"✅ Loaded: {len(prs.slides)} slides")

# ============================================================================
# SLIDE 1: AI Provider Comparison Overview
# ============================================================================
print("\n➕ Adding Slide 1: AI Provider Overview...")

slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
title = slide1.shapes.title
title.text = "🤖 AI Provider Comparison"

# Add subtitle
left = Inches(1)
top = Inches(1.2)
width = Inches(8)
height = Inches(0.5)
subtitle_box = slide1.shapes.add_textbox(left, top, width, height)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Real-world testing: Speed, Cost, and Quality"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

# Add comparison table
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(3.5)

# Table data: 5 rows (header + 4 providers), 5 columns
rows = 5
cols = 5

table = slide1.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)  # Provider
table.columns[1].width = Inches(1.5)  # Speed
table.columns[2].width = Inches(1.5)  # Cost
table.columns[3].width = Inches(1.5)  # Quality
table.columns[4].width = Inches(1.5)  # Best For

# Header row
headers = ['Provider/Model', 'Speed', 'Cost/Query', 'Quality', 'Best For']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # Blue
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White

# Data rows
data = [
    ['GPT-4o-mini', '7.73s ⚡', '$0.0002 💵', '⭐⭐⭐⭐', 'Production'],
    ['Gemini 2.0 Flash', '7.86s ⚡', 'FREE 🎉', '⭐⭐⭐⭐', 'Development'],
    ['Claude Sonnet 4.0', '8.08s', '$0.0036 💰', '⭐⭐⭐⭐⭐', 'Premium'],
    ['LM Studio (Local)', '28.87s 🐌', 'FREE 🎉', '⭐⭐⭐', 'Privacy'],
]

for i, row_data in enumerate(data, start=1):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        
        # Highlight best options
        if i == 1:  # GPT-4o-mini row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 247, 230)  # Light green

print("✅ Slide 1 added")

# ============================================================================
# SLIDE 2: Speed Comparison
# ============================================================================
print("➕ Adding Slide 2: Speed Comparison...")

slide2 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide2.shapes.title
title.text = "⚡ Response Speed Comparison"

# Add speed visualization
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4.5)

textbox = slide2.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Speed bars visualization
speeds = [
    ("GPT-4o-mini", 7.73, "⚡⚡⚡"),
    ("Gemini 2.0 Flash", 7.86, "⚡⚡⚡"),
    ("Claude Sonnet 4.0", 8.08, "⚡⚡"),
    ("LM Studio (Local)", 28.87, "🐌"),
]

p = text_frame.paragraphs[0]
p.text = "Average Response Time (seconds):\n\n"
p.font.size = Pt(16)
p.font.bold = True

for provider, speed, emoji in speeds:
    p = text_frame.add_paragraph()
    bar_length = int((speed / 30) * 40)  # Scale to fit
    bar = "█" * bar_length
    p.text = f"{provider:25} {speed:5.2f}s {emoji}\n{bar}"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)

# Add conclusion
p = text_frame.add_paragraph()
p.text = "\n🏆 Winner: GPT-4o-mini (7.73s average)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)  # Green

print("✅ Slide 2 added")

# ============================================================================
# SLIDE 3: Cost Comparison
# ============================================================================
print("➕ Adding Slide 3: Cost Comparison...")

slide3 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide3.shapes.title
title.text = "💰 Cost Comparison"

# Cost table
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(3)

rows = 5
cols = 4
table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

# Headers
headers = ['Provider', 'Cost/Query', '1K Queries', '1M Queries']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Cost data
cost_data = [
    ['Gemini 2.0 Flash', 'FREE', '$0', '$0 🎉'],
    ['GPT-4o-mini', '$0.0002', '$0.20', '$200 💵'],
    ['Claude Sonnet 4.0', '$0.0036', '$3.60', '$3,600 💰'],
    ['LM Studio', 'FREE', '$0', '$0 🎉'],
]

for i, row_data in enumerate(cost_data, start=1):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        
        # Highlight FREE options
        if 'FREE' in cell_text or '$0' in cell_text:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 247, 230)

# Add key insight
left = Inches(1)
top = Inches(5)
width = Inches(8)
height = Inches(0.8)
insight_box = slide3.shapes.add_textbox(left, top, width, height)
insight_frame = insight_box.text_frame
insight_frame.text = "💡 Key Insight: Claude is 18x more expensive than GPT-4o-mini for similar performance"
insight_para = insight_frame.paragraphs[0]
insight_para.font.size = Pt(14)
insight_para.font.bold = True
insight_para.font.color.rgb = RGBColor(255, 0, 0)

print("✅ Slide 3 added")

# ============================================================================
# SLIDE 4: Quality & Accuracy
# ============================================================================
print("➕ Adding Slide 4: Quality Analysis...")

slide4 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide4.shapes.title
title.text = "📊 Response Quality Analysis"

left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4.5)

textbox = slide4.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Quality analysis
qualities = [
    ("GPT-4o-mini", "⭐⭐⭐⭐", [
        "✓ Clear, well-structured responses",
        "✓ Good balance of empathy and practical advice",
        "✓ Consistent quality across all scenarios",
        "✓ Average 376 tokens (concise)"
    ]),
    ("Gemini 2.0 Flash", "⭐⭐⭐⭐", [
        "✓ Very detailed responses (973 tokens avg)",
        "✓ Comprehensive analysis with multiple perspectives",
        "✓ Excellent markdown formatting",
        "✓ FREE tier available"
    ]),
    ("Claude Sonnet 4.0", "⭐⭐⭐⭐⭐", [
        "✓ Most professional and structured",
        "✓ Concise and actionable (272 tokens)",
        "✓ Superior markdown formatting",
        "✓ Best for high-value interactions"
    ]),
    ("LM Studio", "⭐⭐⭐", [
        "✓ Good detail with tables and formatting",
        "✓ Complete privacy (offline)",
        "✓ No API costs or limits",
        "⚠ 3.7x slower than cloud options"
    ]),
]

for provider, stars, points in qualities:
    p = text_frame.add_paragraph()
    p.text = f"\n{provider} {stars}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(5)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"  {point}"
        p.font.size = Pt(10)
        p.level = 1
        p.space_after = Pt(2)

print("✅ Slide 4 added")

# ============================================================================
# SLIDE 5: Recommendations
# ============================================================================
print("➕ Adding Slide 5: Recommendations...")

slide5 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide5.shapes.title
title.text = "🎯 Recommendations for Socializer"

# Recommendations table
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(2.5)

rows = 5
cols = 3
table = slide5.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(3.0)

# Headers
headers = ['Use Case', 'Recommended Provider', 'Why']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Recommendation data
rec_data = [
    ['Production', 'GPT-4o-mini', 'Fast (7.73s) + Cheap ($0.0002)'],
    ['Development', 'Gemini 2.0 Flash', 'FREE + Good quality'],
    ['Premium Users', 'Claude Sonnet 4.0', 'Best quality + Professional'],
    ['Privacy Mode', 'LM Studio', 'Offline + FREE + Private'],
]

for i, row_data in enumerate(rec_data, start=1):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add cost projection
left = Inches(1)
top = Inches(4.5)
width = Inches(8)
height = Inches(1.5)
cost_box = slide5.shapes.add_textbox(left, top, width, height)
cost_frame = cost_box.text_frame
cost_frame.word_wrap = True

p = cost_frame.paragraphs[0]
p.text = "💵 Estimated Monthly Cost (10,000 users):\n"
p.font.size = Pt(14)
p.font.bold = True

p = cost_frame.add_paragraph()
p.text = "• 7,000 free users → Gemini (FREE)\n• 2,500 standard → GPT-4o-mini ($5/month)\n• 500 premium → Claude ($36/month)"
p.font.size = Pt(12)
p.space_after = Pt(5)

p = cost_frame.add_paragraph()
p.text = "\n🎯 Total: ~$41/month for all AI processing"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)

print("✅ Slide 5 added")

# ============================================================================
# Save presentation
# ============================================================================
output_file = 'Socializer_Presentation_Updated.pptx'
prs.save(output_file)

print(f"\n{'='*60}")
print(f"✅ SUCCESS! Presentation updated")
print(f"{'='*60}")
print(f"📁 Output: {output_file}")
print(f"📊 Total slides: {len(prs.slides)} (added 5 new slides)")
print(f"\nSlides added:")
print("  1. AI Provider Comparison Overview")
print("  2. Speed Comparison")
print("  3. Cost Comparison")
print("  4. Quality & Accuracy Analysis")
print("  5. Recommendations")
print(f"\n💡 Review the presentation and replace 'Socializer_Presentation.pptx' if satisfied")
print(f"{'='*60}")
