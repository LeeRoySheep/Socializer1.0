#!/usr/bin/env python3
"""
Socializer PowerPoint Presentation Creator

Creates a professional PowerPoint presentation for the Socializer application
with database diagrams, system architecture, and placeholders for videos.

Compatible with macOS/PowerPoint for Mac.

Author: Socializer Team
Date: November 12, 2024
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import os

# Modern color scheme
COLORS = {
    'primary': RGBColor(41, 128, 185),      # Blue
    'secondary': RGBColor(52, 152, 219),    # Light Blue
    'accent': RGBColor(46, 204, 113),       # Green
    'dark': RGBColor(44, 62, 80),           # Dark Blue
    'light': RGBColor(236, 240, 241),       # Light Gray
    'white': RGBColor(255, 255, 255),       # White
}


def create_title_slide(prs):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']
    
    # Title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "SOCIALIZER"
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.5), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Social Communication Platform"
    
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(left, top + Inches(3), width, Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"Presentation Date: {datetime.now().strftime('%B %d, %Y')}"
    
    p = date_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_agenda_slide(prs):
    """Create agenda slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    title = slide.shapes.title
    title.text = "📋 Agenda"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    # Content
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    agenda_items = [
        "1. Project Overview",
        "2. System Architecture",
        "3. Database Structure & Relations",
        "4. Key Features",
        "5. Backend API (Swagger Demo)",
        "6. Frontend Demo",
        "7. Security & Encryption",
        "8. Performance & Scalability",
        "9. Future Enhancements"
    ]
    
    for item in agenda_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.level = 0
    
    return slide


def create_overview_slide(prs):
    """Create project overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🎯 Project Overview"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    overview_text = [
        ("What is Socializer?", [
            "AI-powered social communication platform",
            "Real-time chat with AI assistance",
            "Skill tracking and evaluation",
            "Encrypted user data and conversations"
        ]),
        ("Key Technologies", [
            "FastAPI backend (Python)",
            "WebSocket for real-time communication",
            "OpenAI GPT-4 integration",
            "SQLite database with encryption",
            "JWT authentication"
        ])
    ]
    
    for section, items in overview_text:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(24)
        p.font.bold = True
        p.space_before = Pt(20)
        p.level = 0
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_before = Pt(8)
            p.level = 1
    
    return slide


def create_architecture_slide(prs):
    """Create system architecture slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🏗️ System Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    layers = [
        ("Frontend Layer", ["HTML/CSS/JavaScript", "WebSocket Client", "Responsive UI"]),
        ("API Layer", ["FastAPI (45+ endpoints)", "RESTful API", "Swagger Documentation"]),
        ("Business Logic", ["AI Agent (Modular)", "5 Extracted Tools", "3 Handler Classes"]),
        ("Data Layer", ["SQLite Database", "15 Tables", "Full Encryption"]),
        ("External Services", ["OpenAI GPT-4", "Tavily Search", "Email Validation"])
    ]
    
    for layer_name, components in layers:
        p = tf.add_paragraph()
        p.text = layer_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        p.space_before = Pt(12)
        p.level = 0
        
        for component in components:
            p = tf.add_paragraph()
            p.text = f"• {component}"
            p.font.size = Pt(16)
            p.space_before = Pt(6)
            p.level = 1
    
    return slide


def create_database_schema_slide(prs):
    """Create database schema overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🗄️ Database Schema Overview"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Add note about ER diagram
    p = tf.add_paragraph()
    p.text = "📊 Entity-Relationship Diagram"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "[ER Diagram Image Will Be Added Here]"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.space_before = Pt(20)
    
    # Core tables
    p = tf.add_paragraph()
    p.text = "Core Tables (15):"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(30)
    
    tables = [
        "👥 users - User accounts & authentication",
        "💬 chat_rooms - Communication spaces",
        "📝 messages - Chat messages",
        "🎯 skills - Skill definitions",
        "📚 training - User training progress",
        "🎪 life_events - User life events",
        "🔒 user_preferences - Encrypted preferences"
    ]
    
    for table in tables:
        p = tf.add_paragraph()
        p.text = table
        p.font.size = Pt(16)
        p.space_before = Pt(8)
        p.level = 1
    
    return slide


def create_database_details_slide(prs, title_text, table_name, columns, relationships):
    """Create a detailed slide for a specific database table."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = f"📊 {title_text}"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Table name
    p = tf.add_paragraph()
    p.text = f"Table: {table_name}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    p.space_before = Pt(10)
    
    # Columns
    p = tf.add_paragraph()
    p.text = "Columns:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(15)
    
    for col in columns:
        p = tf.add_paragraph()
        p.text = col
        p.font.size = Pt(14)
        p.space_before = Pt(4)
        p.level = 1
    
    # Relationships
    if relationships:
        p = tf.add_paragraph()
        p.text = "Relationships:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_before = Pt(15)
        
        for rel in relationships:
            p = tf.add_paragraph()
            p.text = rel
            p.font.size = Pt(14)
            p.space_before = Pt(4)
            p.level = 1
    
    return slide


def create_features_slide(prs):
    """Create key features slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "✨ Key Features"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    features = [
        ("🤖 AI Integration", [
            "GPT-4 powered conversations",
            "8 specialized AI tools",
            "Context-aware responses",
            "Multi-language support"
        ]),
        ("💬 Real-Time Chat", [
            "WebSocket communication",
            "Public & private rooms",
            "Room invitations",
            "Message history"
        ]),
        ("🔒 Security", [
            "Bcrypt password hashing",
            "Fernet encryption for memory",
            "JWT authentication",
            "User-isolated data"
        ])
    ]
    
    for feature_name, points in features:
        p = tf.add_paragraph()
        p.text = feature_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        p.space_before = Pt(12)
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.space_before = Pt(4)
            p.level = 1
    
    return slide


def create_video_placeholder_slide(prs, title_text, description, video_filename):
    """Create a slide with video placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Video placeholder box
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(4)
    
    video_box = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    video_box.fill.solid()
    video_box.fill.fore_color.rgb = COLORS['light']
    video_box.line.color.rgb = COLORS['primary']
    video_box.line.width = Pt(3)
    
    # Video placeholder text
    text_frame = video_box.text_frame
    text_frame.text = f"🎥 VIDEO: {video_filename}\n\n{description}\n\n[Video will be embedded here]"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # Instruction
    instruction_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    instruction_frame = instruction_box.text_frame
    instruction_frame.text = f"📎 To add video: Insert > Video > Video from File > Select '{video_filename}'"
    p = instruction_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_security_slide(prs):
    """Create security slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🔒 Security & Encryption"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    security_features = [
        ("Password Security", [
            "✅ Bcrypt hashing (100% coverage)",
            "✅ No plain-text passwords stored",
            "✅ Secure password validation"
        ]),
        ("Data Encryption", [
            "✅ Fernet encryption for conversations",
            "✅ User-specific encryption keys",
            "✅ 100% encrypted memory storage"
        ]),
        ("Authentication", [
            "✅ JWT token-based auth",
            "✅ 1-hour token expiration",
            "✅ Token blacklist support"
        ]),
        ("Test Results", [
            "✅ 30/30 users: Passwords hashed",
            "✅ 30/30 users: Have encryption keys",
            "✅ 11/11 users: Memory encrypted"
        ])
    ]
    
    for section, items in security_features:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        p.space_before = Pt(12)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_before = Pt(4)
            p.level = 1
    
    return slide


def create_performance_slide(prs):
    """Create performance slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "⚡ Performance & Quality"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    metrics = [
        ("Code Quality", [
            "✅ 36% code reduction (1,100 lines removed)",
            "✅ Modular architecture",
            "✅ Comprehensive docstrings",
            "✅ Type hints throughout"
        ]),
        ("Testing", [
            "✅ 100% test pass rate",
            "✅ 20+ automated tests",
            "✅ Zero breaking changes",
            "✅ Authentication verified"
        ]),
        ("API Performance", [
            "✅ 45+ endpoints available",
            "✅ OpenAPI documentation",
            "✅ Fast response times",
            "✅ WebSocket real-time updates"
        ])
    ]
    
    for section, items in metrics:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        p.space_before = Pt(12)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.space_before = Pt(6)
            p.level = 1
    
    return slide


def create_future_slide(prs):
    """Create future enhancements slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "🚀 Future Enhancements"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    future_items = [
        "📱 Mobile App Development (iOS/Android)",
        "🌐 Multi-language UI Support",
        "📊 Advanced Analytics Dashboard",
        "🤖 More AI Tools & Capabilities",
        "🔄 Real-time Collaboration Features",
        "📧 Email Notification System",
        "🎨 Customizable UI Themes",
        "🔗 Third-party Integrations (Slack, Discord)",
        "📈 Scalability Improvements",
        "🧠 Advanced Machine Learning Features"
    ]
    
    for item in future_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    return slide


def create_closing_slide(prs):
    """Create closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']
    
    # Thank you message
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    thank_box = slide.shapes.add_textbox(left, top, width, height)
    thank_frame = thank_box.text_frame
    thank_frame.text = "Thank You!"
    
    p = thank_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(2), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions?"
    
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def main():
    """Main function to create the presentation."""
    print("=" * 70)
    print("🎨 Creating Socializer PowerPoint Presentation")
    print("=" * 70)
    print()
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("📄 Creating slides...")
    
    # Create all slides
    slides_created = []
    
    slides_created.append(("Title", create_title_slide(prs)))
    slides_created.append(("Agenda", create_agenda_slide(prs)))
    slides_created.append(("Overview", create_overview_slide(prs)))
    slides_created.append(("Architecture", create_architecture_slide(prs)))
    slides_created.append(("Database Schema", create_database_schema_slide(prs)))
    
    # Database detail slides
    create_database_details_slide(
        prs,
        "Users Table",
        "users",
        [
            "🔑 id (PRIMARY KEY)",
            "👤 username (UNIQUE)",
            "🔒 hashed_password (bcrypt)",
            "🔐 encryption_key (Fernet)",
            "💬 conversation_memory (encrypted JSON)",
            "📧 email",
            "👮 role (User/Admin)",
            "✅ is_active"
        ],
        [
            "→ One-to-Many with chat_rooms (creator)",
            "→ One-to-Many with messages (sender)",
            "→ One-to-Many with life_events",
            "→ Many-to-Many with rooms (via room_members)",
            "→ Many-to-Many with skills (via user_skills)"
        ]
    )
    
    create_database_details_slide(
        prs,
        "Chat Rooms Table",
        "chat_rooms",
        [
            "🔑 id (PRIMARY KEY)",
            "📝 name",
            "👤 creator_id (FOREIGN KEY → users.id)",
            "🕐 created_at",
            "🔓 is_public",
            "🔒 password (optional)",
            "🤖 ai_enabled",
            "✅ is_active"
        ],
        [
            "→ Many-to-One with users (creator)",
            "→ One-to-Many with messages",
            "→ One-to-Many with room_members",
            "→ One-to-Many with room_invites"
        ]
    )
    
    slides_created.append(("Features", create_features_slide(prs)))
    
    # Video placeholder slides
    slides_created.append(("Backend Video", create_video_placeholder_slide(
        prs,
        "🔧 Backend API Demo (Swagger UI)",
        "Demonstration of API endpoints using Swagger UI\n• Authentication\n• API testing\n• Real-time responses",
        "swagger_demo.mp4"
    )))
    
    slides_created.append(("Frontend Video", create_video_placeholder_slide(
        prs,
        "🎨 Frontend Demo",
        "User interface walkthrough\n• Chat functionality\n• Room management\n• AI interactions",
        "frontend_demo.mp4"
    )))
    
    slides_created.append(("Security", create_security_slide(prs)))
    slides_created.append(("Performance", create_performance_slide(prs)))
    slides_created.append(("Future", create_future_slide(prs)))
    slides_created.append(("Closing", create_closing_slide(prs)))
    
    print(f"✅ Created {len(slides_created)} slides")
    print()
    
    # Save presentation
    filename = "Socializer_Presentation.pptx"
    prs.save(filename)
    
    print("=" * 70)
    print(f"✅ Presentation saved: {filename}")
    print("=" * 70)
    print()
    print("📊 Next Steps:")
    print("1. Open the presentation in PowerPoint")
    print("2. Review all slides")
    print("3. Add ER diagram image to 'Database Schema Overview' slide")
    print("4. Record and embed videos:")
    print("   - swagger_demo.mp4 (Backend API demo)")
    print("   - frontend_demo.mp4 (Frontend walkthrough)")
    print("5. Customize as needed")
    print()
    print("🎥 Video Recording Instructions:")
    print("- Use QuickTime Player (Cmd+Shift+5) for screen recording")
    print("- Save videos in the same folder as this presentation")
    print("- In PowerPoint: Insert > Video > Video from File")
    print()
    
    return filename


if __name__ == "__main__":
    try:
        filename = main()
        print(f"🎉 Success! Open '{filename}' in PowerPoint for Mac")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
